import argparse
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
import cohere
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
import os

# Add Cohere key
COHERE_API_KEY = os.getenv("COHERE_API_KEY")
co = cohere.Client(COHERE_API_KEY)

def get_inputs():
    parser = argparse.ArgumentParser(description="Automate prospect company research.")
    parser.add_argument("--website", required=True, help="Company website URL")
    parser.add_argument("--emails", required=True, help="Path to a .txt file with email addresses")
    args = parser.parse_args()
    return args.website, args.emails

def scrape_website(url):
    response = requests.get(url)
    soup = BeautifulSoup(response.text, "html.parser")

    # Try to extract meaningful content
    main_content = soup.find("div", {"id": "mw-content-text"})  # Specific to Wikipedia
    if main_content:
        return main_content.text.strip()

    # Fallback if no specific content is found
    return "Content not found on the provided webpage."

def generate_informative_article(text):
    """Generates a detailed informative article based on the provided text."""
    response = co.generate(
        model="command-xlarge-nightly",  # Replace with the appropriate model
        prompt=(
            f"Write a detailed and informative article based on the text below. "
            f"If the base text is insufficient, use general knowledge to complete the content.\n\n"
            f"Base text:\n{text}\n\n"
            f"Informative article:"
        ),
        max_tokens=2000,  # Adjust based on your requirements
        temperature=0.7
    )
    return response.generations[0].text.strip()

def summarize_content(content):
    """Resuma o conteúdo para criar tópicos mais concisos."""
    response = co.generate(
        model="command-xlarge-nightly",
        prompt=(
            f"Resuma o seguinte texto em tópicos curtos e claros para uma apresentação:"
            f"\n\n{content}\n\nResumo conciso:"
        ),
        max_tokens=500,
        temperature=0.7
    )
    return response.generations[0].text.strip()

def create_presentation(content):
    prs = Presentation()
    
    # Slide inicial
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Company Overview"
    title_slide.placeholders[1].text = "The Evolution of the Electric Guitar"

    # Dividir conteúdo em seções
    sections = content.split("\n\n")  # Dividir em parágrafos
    for i, section in enumerate(sections):
        if len(section.strip()) > 0:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Título e conteúdo
            slide.shapes.title.text = f"Seção {i + 1}"  # Adicionar título dinâmico
            bullets = section.split(". ")  # Dividir o texto em frases
            content_box = slide.placeholders[1]
            
            for bullet in bullets:
                if bullet.strip():  # Ignorar textos vazios
                    content_box.text += f"- {bullet.strip()}.\n"  # Adicionar como bullet point

    prs.save("summary.pptx")

def send_email(recipients, subject, body, attachment):
    sender_email = os.getenv("EMAIL_ADDRESS")
    password = os.getenv("EMAIL_PASSWORD")
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['Subject'] = subject
    msg.attach(MIMEText(body, 'plain'))
    with open(attachment, "rb") as file:
        part = MIMEBase("application", "octet-stream")
        part.set_payload(file.read())
        encoders.encode_base64(part)
        part.add_header("Content-Disposition", f"attachment; filename={attachment}")
        msg.attach(part)
    with smtplib.SMTP('smtp.gmail.com', 587) as server:
        server.starttls()
        server.login(sender_email, password)
        for recipient in recipients:
            server.sendmail(sender_email, recipient, msg.as_string())

if __name__ == "__main__":
    website, email_file = get_inputs()
    with open(email_file, "r") as file:
        emails = file.read().splitlines()
    
    # Extract text from the website
    scraped_text = scrape_website(website)
    
    # Generate an informative article
    article = generate_informative_article(scraped_text)
    
    # Create a presentation with the article content
    create_presentation(article)
    
    # Send the article via email
    send_email(emails, "Company Informative Article", article, "summary.pptx")