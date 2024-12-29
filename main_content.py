import argparse
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
import openai
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
import os

# Configure the API key
openai.api_key = os.getenv("OPENAI_API_KEY")  # Retrieve the API key from an environment variable

def get_inputs():
    parser = argparse.ArgumentParser(description="Automate prospect company research.")
    parser.add_argument("--website", required=True, help="Company website URL")
    parser.add_argument("--emails", required=True, help="Path to a .txt file with email addresses")
    args = parser.parse_args()
    return args.website, args.emails

def scrape_website(url):
    response = requests.get(url)
    soup = BeautifulSoup(response.text, "html.parser")
    about_us = soup.find("section", {"id": "about"}) or soup.find("div", text="About Us")
    return about_us.text.strip() if about_us else "About Us not found."

def generate_informative_article(text):
    """Generates a detailed informative article based on the provided text."""
    response = openai.ChatCompletion.create(
    model="gpt-3.5-turbo",
    messages=[
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": (
            f"Write a detailed and informative article based on the text below. "
            f"Make sure to include an introduction, subtitles for each section, and a conclusion.\n\n"
            f"Base text:\n{text}\n\n"
            f"Informative article:"
        )}
    ],
    max_tokens=1000
)

    return response['choices'][0]['message']['content'].strip()


def create_presentation(content):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Company Overview"
    slide.placeholders[1].text = content
    prs.save("summary.pptx")

def send_email(recipients, subject, body, attachment):
    sender_email = os.getenv("EMAIL_ADDRESS")
    password = os.getenv("EMAIL_PASSWORD")
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['Subject'] = subject
    msg.attach(MIMEText(body, 'plain'))
    with open(attachment, "rb") as file:
        part = MIMEBase("application", "octet-stream")
        part.set_payload(file.read())
        encoders.encode_base64(part)
        part.add_header("Content-Disposition", f"attachment; filename={attachment}")
        msg.attach(part)
    with smtplib.SMTP('smtp.gmail.com', 587) as server:
        server.starttls()
        server.login(sender_email, password)
        for recipient in recipients:
            server.sendmail(sender_email, recipient, msg.as_string())

if __name__ == "__main__":
    website, email_file = get_inputs()
    with open(email_file, "r") as file:
        emails = file.read().splitlines()
    
    # Extract text from the website
    scraped_text = scrape_website(website)
    
    # Generate an informative article
    article = generate_informative_article(scraped_text)
    
    # Create a presentation with the article content
    create_presentation(article)
    
    # Send the article via email
    send_email(emails, "Company Informative Article", article, "summary.pptx")